from pptx import Presentation
import base64
import pandas as pd

# Extracts the data from the PPT file
def load(file):
    ppt = Presentation(file)
    ppt_data = {}
    for idx, slide in enumerate(ppt.slides):
        ppt_data[idx] = extract_slide(slide)
    return ppt_data

# Extracts data from a PPT slide
def extract_slide(slide):
    slide_data = []

    for shape in slide.shapes:
        if shape.shape_type == 3:
            print("chart")
        # Extracts data for each shape_type based on the switch condition
        slide_data.append(switch_type(shape.shape_type)(shape))

    return slide_data

def extract_shape(shape):
    return shape.shape_type

def binary_to_decoded_bytes(bytes):
    return base64.b64encode(bytes).decode("utf-8")

def extract_picture(shape):
    image_data = binary_to_decoded_bytes(shape.image.blob)
    data = {
        "data" : image_data,
        "shape_type" : shape.shape_type
    }
    return data

def extract_text(shape):
    text_data = shape.text
    data = {
        "data" : text_data,
        "shape_type" : shape.shape_type
    }
    return data

def extract_placeholder(shape):
    if shape.has_text_frame:
        return extract_text(shape)
    elif shape.has_table:
        return extract_table(shape)
    elif shape.has_chart:
        return extract_slide(shape)
    data = {
        "data" : "Unknown Placeholder data",
        "shape_type" : shape.shape_type
    }
    return data

def extract_table(shape):
    rows_data = []
    for row in shape.table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text)
        rows_data.append(row_data)
    df = pd.DataFrame(rows_data)

    # Code for having first row as a column
    # df.columns = df.iloc[0]
    # df = df[1:]

    table_data = "The following is a table data with default columns : " + "\n"
    for index, row in df.iterrows():
        for index, item in row.items():
            table_data += f"{index} : {item}"
            table_data += ", "
        table_data += "\n"
    
    data = {
        "data" : table_data,
        "shape_type" : shape.shape_type
    }
    return data

# Switches to different data extractors based on the shape type. Returns the function object which can be run in the original function
def switch_type(shape_type):
    switch = {
        -2 : extract_shape, # "mixed",
        1 : extract_shape, # "autoshape",
        2 : extract_shape, # "callout",
        3 : extract_shape, # "chart",
        4 : extract_shape, # "comment",
        5 : extract_shape, # "freeform",
        6 : extract_slide, # "group"
        7 : extract_shape, # "embedded_ole_object",
        8 : extract_shape, # "form_control",
        9 : extract_shape, # "line",
        10 : extract_shape, # "linked_ole_object",
        11 : extract_shape, # "linked_picture",
        12 : extract_shape, # "ole_control_object",
        13 : extract_picture, # "picture",
        14 : extract_placeholder, # "placeholder",
        16 : extract_shape, # "media",
        17 : extract_text, # "textbox",
        18 : extract_shape, # "script_anchor",
        19 : extract_table, # "table",
        20 : extract_shape, # "canvas",
        21 : extract_shape, # "diagram",
        22 : extract_shape, # "ink",
        23 : extract_shape, # "ink_comment",
        24 : extract_shape, # "igx_graphic",
        25 : extract_shape, # "",
        26 : extract_shape, # "web_video"
    }
    return switch.get(shape_type)